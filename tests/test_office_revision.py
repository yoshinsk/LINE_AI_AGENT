r"""<PROJECT_ROOT>\tests\test_office_revision.py

構造化したOffice編集計画から、実際に修正済みDOCX/XLSXを返却できることを検証します。
"""

from __future__ import annotations

from pathlib import Path
import tempfile
import unittest

import fitz
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font
from pptx import Presentation

from line_ai_agent.office_revision import OfficeRevisionError, apply_office_revision_plan


class OfficeRevisionTest(unittest.TestCase):
    """Word/Excelの実ファイルに安全な範囲で編集計画を反映する経路を確認します。"""

    def test_revises_docx_text_across_runs(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "motivation.docx"
            document = Document()
            paragraph = document.add_paragraph()
            paragraph.add_run("交換")
            paragraph.add_run("留学を志望します。")
            document.save(source)

            results = apply_office_revision_plan(
                {
                    "summary": "修正済みWord文書を作成しました。",
                    "files": [
                        {
                            "source_file": source.name,
                            "edits": [{"kind": "word_text", "find": "交換留学", "replacement": "海外留学"}],
                        }
                    ],
                },
                (source,),
                root / "output",
            )

            self.assertEqual((root / "output" / "motivation-revised.docx",), results)
            self.assertEqual("海外留学を志望します。", Document(results[0]).paragraphs[0].text)

    def test_revises_xlsx_cell_without_changing_formula_or_style(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "schools.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "志望校"
            sheet["A1"] = "第一希望"
            sheet["A1"].font = Font(bold=True)
            sheet["B1"] = "=SUM(1,2)"
            workbook.save(source)

            results = apply_office_revision_plan(
                {
                    "summary": "修正済みExcel文書を作成しました。",
                    "files": [
                        {
                            "source_file": source.name,
                            "edits": [
                                {
                                    "kind": "excel_cell",
                                    "sheet": "志望校",
                                    "cell": "A1",
                                    "expected": "第一希望",
                                    "replacement": "第一志望",
                                }
                            ],
                        }
                    ],
                },
                (source,),
                root / "output",
            )

            revised = load_workbook(results[0], data_only=False)["志望校"]
            self.assertEqual("第一志望", revised["A1"].value)
            self.assertTrue(revised["A1"].font.bold)
            self.assertEqual("=SUM(1,2)", revised["B1"].value)

    def test_rejects_plan_that_does_not_modify_the_source(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.docx"
            document = Document()
            document.add_paragraph("原文")
            document.save(source)

            with self.assertRaises(OfficeRevisionError):
                apply_office_revision_plan(
                    {
                        "summary": "修正済みWord文書を作成しました。",
                        "files": [
                            {
                                "source_file": source.name,
                                "edits": [{"kind": "word_text", "find": "存在しない文章", "replacement": "置換"}],
                            }
                        ],
                    },
                    (source,),
                    root / "output",
                )

    def test_revises_pptx_text(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "presentation.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            shape = slide.shapes.add_textbox(0, 0, 5_000_000, 500_000)
            shape.text_frame.paragraphs[0].add_run().text = "旧タイトル"
            presentation.save(source)

            results = apply_office_revision_plan(
                {
                    "summary": "修正済みPowerPointを作成しました。",
                    "files": [
                        {
                            "source_file": source.name,
                            "edits": [{"kind": "pptx_text", "find": "旧タイトル", "replacement": "新タイトル"}],
                        }
                    ],
                },
                (source,),
                root / "output",
            )

            revised = Presentation(results[0])
            texts = [shape.text for shape in revised.slides[0].shapes if getattr(shape, "has_text_frame", False)]
            self.assertIn("新タイトル", texts)

    def test_revises_searchable_pdf_text(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "report.pdf"
            document = fitz.open()
            page = document.new_page()
            page.insert_text((72, 72), "Original wording")
            document.save(source)
            document.close()

            results = apply_office_revision_plan(
                {
                    "summary": "修正済みPDFを作成しました。",
                    "files": [
                        {
                            "source_file": source.name,
                            "edits": [{"kind": "pdf_text", "find": "Original", "replacement": "Revised"}],
                        }
                    ],
                },
                (source,),
                root / "output",
            )

            revised = fitz.open(results[0])
            try:
                page_text = revised[0].get_text("text")
            finally:
                revised.close()
            self.assertIn("Revised", page_text)
            self.assertNotIn("Original", page_text)


if __name__ == "__main__":
    unittest.main()
