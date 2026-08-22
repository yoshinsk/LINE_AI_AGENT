r"""<PROJECT_ROOT>\src\line_ai_agent\office_revision.py

Codexが返した構造化したOffice編集計画をDOCX/XLSXへ反映し、LINE送信用の修正済みファイルを生成します。
"""

from __future__ import annotations

from collections.abc import Iterable
from pathlib import Path
import os
import posixpath
from typing import Any
import xml.etree.ElementTree as ElementTree
import zipfile

import fitz
from docx import Document
from docx.document import Document as WordDocument
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class OfficeRevisionError(ValueError):
    """編集計画を安全にOffice文書へ反映できない場合に送出します。"""


REVISION_DOCUMENT_SUFFIXES = {".docx", ".xlsx", ".pptx", ".pdf"}
SPREADSHEET_NS = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
OFFICE_RELATIONSHIPS_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PACKAGE_RELATIONSHIPS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def apply_office_revision_plan(plan: dict[str, Any], sources: tuple[Path, ...], output_dir: Path) -> tuple[Path, ...]:
    """編集計画の各Word、Excel、PowerPoint、PDFを修正し、実際に編集できた成果物だけを返します。"""
    source_by_name = {path.name: path for path in sources if path.suffix.lower() in REVISION_DOCUMENT_SUFFIXES}
    file_plans = plan.get("files")
    if not isinstance(file_plans, list) or not file_plans:
        raise OfficeRevisionError("編集計画にfilesがありません。")

    output_dir.mkdir(parents=True, exist_ok=True)
    results: list[Path] = []
    errors: list[str] = []
    planned_names: set[str] = set()
    for file_plan in file_plans:
        if not isinstance(file_plan, dict):
            errors.append("編集対象の指定が不正です。")
            continue
        source_name = Path(str(file_plan.get("source_file") or "")).name
        source = source_by_name.get(source_name)
        if source is None:
            errors.append(f"対象ファイルが見つかりません: {source_name or '未指定'}")
            continue
        if source.name in planned_names:
            errors.append(f"同じファイルが複数回指定されています: {source.name}")
            continue
        planned_names.add(source.name)
        try:
            result = _revise_file(source, file_plan, output_dir)
        except OfficeRevisionError as exc:
            errors.append(f"{source.name}: {exc}")
            continue
        results.append(result)

    missing = sorted(set(source_by_name) - planned_names)
    if missing:
        errors.append("編集計画にない添付ファイル: " + ", ".join(missing))
    if errors or len(results) != len(source_by_name):
        detail = " / ".join(errors) or "修正済みファイルが不足しています。"
        raise OfficeRevisionError(detail)
    return tuple(results)


def _revise_file(source: Path, file_plan: dict[str, Any], output_dir: Path) -> Path:
    """拡張子ごとの編集を実行し、変更件数が0の単なるコピーを返さないようにします。"""
    edits = file_plan.get("edits")
    if not isinstance(edits, list) or not edits:
        raise OfficeRevisionError("実際に反映するeditsがありません。")
    destination = output_dir / f"{source.stem}-revised{source.suffix.lower()}"
    suffix = source.suffix.lower()
    if suffix == ".docx":
        change_count = _revise_docx(source, destination, edits)
    elif suffix == ".xlsx":
        change_count = _revise_xlsx(source, destination, edits)
    elif suffix == ".pptx":
        change_count = _revise_pptx(source, destination, edits)
    elif suffix == ".pdf":
        change_count = _revise_pdf(source, destination, edits)
    else:
        raise OfficeRevisionError("未対応のOffice形式です。")
    if change_count <= 0:
        destination.unlink(missing_ok=True)
        raise OfficeRevisionError("原文と一致する編集箇所がなく、修正を反映できませんでした。")
    return destination.resolve()


def _revise_docx(source: Path, destination: Path, edits: list[Any]) -> int:
    """Word本文・表・ヘッダー等の段落に対して、計画内の完全一致置換を反映します。"""
    replacements = [_word_replacement(item) for item in edits]
    replacements = [item for item in replacements if item is not None]
    if not replacements:
        raise OfficeRevisionError("Word用の文字列置換がありません。")

    document = Document(source)
    change_count = 0
    for paragraph in _iter_document_paragraphs(document):
        for find, replacement in replacements:
            change_count += _replace_paragraph_text(paragraph, find, replacement)
    if change_count:
        document.save(destination)
    return change_count


def _revise_xlsx(source: Path, destination: Path, edits: list[Any]) -> int:
    """Excelの対象セルXMLだけを置換し、Open XML内の数式・検証・拡張要素をそのまま保持します。"""
    changes = [change for item in edits if (change := _excel_cell_change(item)) is not None]
    if not changes:
        raise OfficeRevisionError("Excel用のセル置換がありません。")

    try:
        with zipfile.ZipFile(source) as archive:
            members = {info.filename: archive.read(info.filename) for info in archive.infolist()}
            sheet_members = _xlsx_sheet_members(members)
            shared_root, shared_values, shared_member, has_shared_strings = _xlsx_shared_strings(members)
            edited_sheets: dict[str, ElementTree.Element] = {}
            change_count = 0
            for sheet_name, coordinate, expected, replacement in changes:
                member_name = sheet_members.get(sheet_name)
                if member_name is None:
                    continue
                root = edited_sheets.get(member_name)
                if root is None:
                    raw = members.get(member_name)
                    if raw is None:
                        continue
                    root = ElementTree.fromstring(raw)
                    edited_sheets[member_name] = root
                cell = root.find(f".//{{{SPREADSHEET_NS}}}c[@r='{coordinate}']")
                if cell is None or cell.find(f"{{{SPREADSHEET_NS}}}f") is not None:
                    continue
                actual = _xlsx_cell_value(cell, shared_values)
                if actual != expected or actual == replacement:
                    continue
                if has_shared_strings:
                    shared_index = _append_xlsx_shared_string(shared_root, shared_values, replacement)
                    _set_xlsx_cell_shared_string(cell, shared_index)
                else:
                    _set_xlsx_cell_inline_string(cell, replacement)
                change_count += 1
            if not change_count:
                return 0
            for member_name, root in edited_sheets.items():
                members[member_name] = ElementTree.tostring(root, encoding="utf-8", xml_declaration=True)
            if has_shared_strings:
                members[shared_member] = ElementTree.tostring(shared_root, encoding="utf-8", xml_declaration=True)
            _write_xlsx_archive(source, destination, members)
            return change_count
    except (ElementTree.ParseError, KeyError, OSError, ValueError, zipfile.BadZipFile) as exc:
        raise OfficeRevisionError(f"XLSXを安全に更新できませんでした: {exc}") from exc


def _xlsx_sheet_members(members: dict[str, bytes]) -> dict[str, str]:
    """workbook.xmlと関連定義から、シート名ごとのワークシートXMLを解決します。"""
    workbook_raw = members.get("xl/workbook.xml")
    relationships_raw = members.get("xl/_rels/workbook.xml.rels")
    if workbook_raw is None or relationships_raw is None:
        raise ValueError("XLSXのworkbook関連定義がありません。")
    workbook = ElementTree.fromstring(workbook_raw)
    relationships = ElementTree.fromstring(relationships_raw)
    targets = {
        item.attrib.get("Id", ""): item.attrib.get("Target", "")
        for item in relationships.findall(f"{{{PACKAGE_RELATIONSHIPS_NS}}}Relationship")
    }
    sheets: dict[str, str] = {}
    for item in workbook.findall(f".//{{{SPREADSHEET_NS}}}sheet"):
        sheet_name = item.attrib.get("name", "")
        relation_id = item.attrib.get(f"{{{OFFICE_RELATIONSHIPS_NS}}}id", "")
        target = targets.get(relation_id, "")
        member_name = _xlsx_member_name(target)
        if sheet_name and member_name in members:
            sheets[sheet_name] = member_name
    return sheets


def _xlsx_member_name(target: str) -> str:
    """workbook関連の相対パスを、XLSXコンテナ内のxl配下パスへ正規化します。"""
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join("xl", target))


def _xlsx_shared_strings(members: dict[str, bytes]) -> tuple[ElementTree.Element, list[str], str, bool]:
    """共有文字列を読み、元の文字列を変更せず新規文字列を追加できる状態にします。"""
    member_name = "xl/sharedStrings.xml"
    raw = members.get(member_name)
    if raw is None:
        root = ElementTree.Element(f"{{{SPREADSHEET_NS}}}sst", {"count": "0", "uniqueCount": "0"})
        return root, [], member_name, False
    root = ElementTree.fromstring(raw)
    values = ["".join(text.text or "" for text in item.iter(f"{{{SPREADSHEET_NS}}}t")) for item in root.findall(f"{{{SPREADSHEET_NS}}}si")]
    return root, values, member_name, True


def _xlsx_cell_value(cell: ElementTree.Element, shared_values: list[str]) -> str:
    """対象セルの表示文字列を取得し、編集計画のexpectedと厳密比較できるようにします。"""
    cell_type = cell.attrib.get("t", "")
    if cell_type == "s":
        value = cell.findtext(f"{{{SPREADSHEET_NS}}}v", default="")
        try:
            return shared_values[int(value)]
        except (IndexError, ValueError):
            return ""
    if cell_type == "inlineStr":
        return "".join(text.text or "" for text in cell.iter(f"{{{SPREADSHEET_NS}}}t"))
    return cell.findtext(f"{{{SPREADSHEET_NS}}}v", default="")


def _append_xlsx_shared_string(root: ElementTree.Element, values: list[str], value: str) -> int:
    """既存参照へ影響しない新しい共有文字列を追加し、そのインデックスを返します。"""
    index = len(values)
    item = ElementTree.SubElement(root, f"{{{SPREADSHEET_NS}}}si")
    text = ElementTree.SubElement(item, f"{{{SPREADSHEET_NS}}}t")
    if value.startswith(" ") or value.endswith(" "):
        text.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    text.text = value
    values.append(value)
    root.set("count", str(max(0, int(root.attrib.get("count", "0"))) + 1))
    root.set("uniqueCount", str(len(values)))
    return index


def _set_xlsx_cell_shared_string(cell: ElementTree.Element, shared_index: int) -> None:
    """セルの値要素だけを共有文字列参照へ置き換え、style・検証・コメントなどの属性を残します。"""
    cell.set("t", "s")
    for child in list(cell):
        cell.remove(child)
    value = ElementTree.SubElement(cell, f"{{{SPREADSHEET_NS}}}v")
    value.text = str(shared_index)


def _set_xlsx_cell_inline_string(cell: ElementTree.Element, replacement: str) -> None:
    """共有文字列テーブルがないブックでは、対象セルだけをinline文字列へ変更して関係定義を増やさないようにします。"""
    cell.set("t", "inlineStr")
    for child in list(cell):
        cell.remove(child)
    inline = ElementTree.SubElement(cell, f"{{{SPREADSHEET_NS}}}is")
    text = ElementTree.SubElement(inline, f"{{{SPREADSHEET_NS}}}t")
    if replacement.startswith(" ") or replacement.endswith(" "):
        text.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    text.text = replacement


def _write_xlsx_archive(source: Path, destination: Path, members: dict[str, bytes]) -> None:
    """変更対象XML以外は元のZipInfoとバイト列を維持し、Office拡張要素を失わずに保存します。"""
    with zipfile.ZipFile(source) as input_archive, zipfile.ZipFile(destination, "w") as output_archive:
        for info in input_archive.infolist():
            output_archive.writestr(info, members[info.filename])


def _revise_pptx(source: Path, destination: Path, edits: list[Any]) -> int:
    """PowerPointのテキストボックスと表セルの文字列を、既存のテーマ・配置を保って置換します。"""
    replacements = [_pptx_replacement(item) for item in edits]
    replacements = [item for item in replacements if item is not None]
    if not replacements:
        raise OfficeRevisionError("PowerPoint用の文字列置換がありません。")

    presentation = Presentation(source)
    change_count = 0
    for slide in presentation.slides:
        for text_frame in _iter_pptx_text_frames(slide.shapes):
            for paragraph in text_frame.paragraphs:
                for find, replacement in replacements:
                    change_count += _replace_runs_text(list(paragraph.runs), find, replacement)
    if change_count:
        presentation.save(destination)
    return change_count


def _revise_pdf(source: Path, destination: Path, edits: list[Any]) -> int:
    """PDFの検索可能な文字列を赤塗り置換し、元ページのレイアウトを保った修正済みPDFを保存します。"""
    replacements = [_pdf_replacement(item) for item in edits]
    replacements = [item for item in replacements if item is not None]
    if not replacements:
        raise OfficeRevisionError("PDF用の文字列置換がありません。")

    document = fitz.open(source)
    try:
        if document.needs_pass:
            raise OfficeRevisionError("パスワード保護されたPDFは編集できません。")
        change_count = 0
        for page in document:
            pending: list[tuple[fitz.Rect, str]] = []
            for find, replacement in replacements:
                for rectangle in page.search_for(find):
                    page.add_redact_annot(rectangle, fill=(1, 1, 1))
                    pending.append((rectangle, replacement))
            if not pending:
                continue
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
            for rectangle, replacement in pending:
                _insert_pdf_replacement(page, rectangle, replacement)
            change_count += len(pending)
        if change_count:
            document.save(destination, garbage=4, deflate=True)
        return change_count
    finally:
        document.close()


def _word_replacement(item: Any) -> tuple[str, str] | None:
    """Word向け編集だけを受け付け、空文字・同一文字列の置換を捨てます。"""
    if not isinstance(item, dict) or item.get("kind") != "word_text":
        return None
    find = item.get("find")
    replacement = item.get("replacement")
    if not isinstance(find, str) or not isinstance(replacement, str) or not find or find == replacement:
        return None
    return find, replacement


def _pptx_replacement(item: Any) -> tuple[str, str] | None:
    """PowerPoint向け編集だけを受け付け、本文と一致しない曖昧な置換を捨てます。"""
    if not isinstance(item, dict) or item.get("kind") != "pptx_text":
        return None
    find = item.get("find")
    replacement = item.get("replacement")
    if not isinstance(find, str) or not isinstance(replacement, str) or not find or find == replacement:
        return None
    return find, replacement


def _pdf_replacement(item: Any) -> tuple[str, str] | None:
    """PDF向け編集だけを受け付け、検索不能な空文字列の置換を防ぎます。"""
    if not isinstance(item, dict) or item.get("kind") != "pdf_text":
        return None
    find = item.get("find")
    replacement = item.get("replacement")
    if not isinstance(find, str) or not isinstance(replacement, str) or not find or find == replacement:
        return None
    return find, replacement


def _excel_cell_change(item: Any) -> tuple[str, str, str, str] | None:
    """Excel向け編集を検証し、意図しないセル更新を避けるため期待値も必須にします。"""
    if not isinstance(item, dict) or item.get("kind") != "excel_cell":
        return None
    sheet = item.get("sheet")
    cell = item.get("cell")
    expected = item.get("expected")
    replacement = item.get("replacement")
    if not all(isinstance(value, str) and value for value in (sheet, cell, expected, replacement)):
        return None
    return sheet, cell, expected, replacement


def _iter_document_paragraphs(document: WordDocument) -> Iterable[Paragraph]:
    """本文・表・ヘッダー・フッターを含む編集可能な段落を重複なく列挙します。"""
    seen: set[int] = set()
    for paragraph in document.paragraphs:
        yield from _yield_once(paragraph, seen)
    for table in document.tables:
        yield from _iter_table_paragraphs(table, seen)
    for section in document.sections:
        for part in (section.header, section.footer, section.first_page_header, section.first_page_footer):
            for paragraph in part.paragraphs:
                yield from _yield_once(paragraph, seen)
            for table in part.tables:
                yield from _iter_table_paragraphs(table, seen)


def _iter_table_paragraphs(table: Table, seen: set[int]) -> Iterable[Paragraph]:
    """入れ子の表を含めてセル段落を走査します。"""
    for row in table.rows:
        for cell in row.cells:
            yield from _iter_cell_paragraphs(cell, seen)


def _iter_cell_paragraphs(cell: _Cell, seen: set[int]) -> Iterable[Paragraph]:
    """セル本文とセル内の入れ子表を走査します。"""
    for paragraph in cell.paragraphs:
        yield from _yield_once(paragraph, seen)
    for table in cell.tables:
        yield from _iter_table_paragraphs(table, seen)


def _yield_once(paragraph: Paragraph, seen: set[int]) -> Iterable[Paragraph]:
    """同じXML段落を二重に編集しないよう、内部要素のIDで重複を除きます。"""
    key = id(paragraph._p)
    if key not in seen:
        seen.add(key)
        yield paragraph


def _iter_pptx_text_frames(shapes: Any) -> Iterable[Any]:
    """グループ化された図形と表セルを含め、PowerPoint内の編集可能なテキスト枠を列挙します。"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_pptx_text_frames(shape.shapes)
        if getattr(shape, "has_text_frame", False):
            yield shape.text_frame
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


def _replace_paragraph_text(paragraph: Paragraph, find: str, replacement: str) -> int:
    """段落内のrunをまたぐ完全一致文字列を後方から置換し、既存の書式をできる限り維持します。"""
    return _replace_runs_text(list(paragraph.runs), find, replacement)


def _replace_runs_text(runs: list[Any], find: str, replacement: str) -> int:
    """WordとPowerPointのrun列を対象に、書式を保持しながら完全一致文字列を後方から置換します。"""
    if not runs:
        return 0
    original = "".join(run.text for run in runs)
    starts: list[int] = []
    cursor = 0
    while True:
        index = original.find(find, cursor)
        if index < 0:
            break
        starts.append(index)
        cursor = index + len(find)
    for start in reversed(starts):
        _replace_run_span(runs, start, start + len(find), replacement)
    return len(starts)


def _replace_run_span(runs: list[Any], start: int, end: int, replacement: str) -> None:
    """複数runにまたがる範囲を置換し、置換前後の文字列を元runへ残します。"""
    first_index, first_offset = _locate_run_offset(runs, start)
    last_index, last_offset = _locate_run_offset(runs, end - 1)
    if first_index == last_index:
        value = runs[first_index].text
        runs[first_index].text = value[:first_offset] + replacement + value[last_offset + 1 :]
        return
    first_value = runs[first_index].text
    last_value = runs[last_index].text
    runs[first_index].text = first_value[:first_offset] + replacement
    for index in range(first_index + 1, last_index):
        runs[index].text = ""
    runs[last_index].text = last_value[last_offset + 1 :]


def _locate_run_offset(runs: list[Any], position: int) -> tuple[int, int]:
    """連結済み段落文字列の位置を、対応するrunとrun内オフセットへ変換します。"""
    offset = 0
    for index, run in enumerate(runs):
        next_offset = offset + len(run.text)
        if position < next_offset:
            return index, position - offset
        offset = next_offset
    raise OfficeRevisionError("Word段落の置換位置を特定できませんでした。")


def _insert_pdf_replacement(page: fitz.Page, rectangle: fitz.Rect, replacement: str) -> None:
    """元の文字領域へ日本語対応フォントで置換文字列を描画し、収まらない場合は失敗として扱います。"""
    font_path = _pdf_japanese_font_path()
    font_size = max(6.0, min(18.0, rectangle.height * 0.82))
    options: dict[str, Any] = {"fontname": "helv", "color": (0, 0, 0), "overlay": True}
    if font_path is not None:
        options = {"fontname": "line_agent_japanese", "fontfile": str(font_path), "color": (0, 0, 0), "overlay": True}
    for size in (font_size, font_size * 0.8, font_size * 0.65, font_size * 0.5, 4.0):
        result = page.insert_textbox(rectangle, replacement, fontsize=max(4.0, size), **options)
        if result >= 0:
            return
    raise OfficeRevisionError("PDFの置換文字列が元の領域に収まりません。")


def _pdf_japanese_font_path() -> Path | None:
    """Windowsワーカー上では游ゴシックまたはメイリオを埋め込み、日本語PDFの文字化けを防ぎます。"""
    font_dir = Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts"
    for name in ("YuGothM.ttc", "meiryo.ttc", "msgothic.ttc"):
        candidate = font_dir / name
        if candidate.is_file():
            return candidate
    return None
